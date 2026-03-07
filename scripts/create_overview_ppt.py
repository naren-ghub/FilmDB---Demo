"""FilmDB Demo — Overview Presentation Generator"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color Palette ──────────────────────────────────────────────────────────
BG_DARK     = RGBColor(0x0F, 0x0F, 0x1A)   # Deep navy
BG_CARD     = RGBColor(0x1A, 0x1A, 0x2E)   # Card background
ACCENT_GOLD = RGBColor(0xFF, 0xC1, 0x07)   # Gold accent
ACCENT_BLUE = RGBColor(0x4E, 0x9A, 0xF5)   # Blue accent
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xAA)   # Teal/cyan
ACCENT_PURPLE = RGBColor(0xBB, 0x86, 0xFC) # Purple
ACCENT_ORANGE = RGBColor(0xFF, 0x8A, 0x65) # Orange
ACCENT_RED  = RGBColor(0xEF, 0x53, 0x50)   # Red
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xB0, 0xB0, 0xC0)
DIM_GRAY    = RGBColor(0x80, 0x80, 0x99)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def _set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if alpha is not None:
        from lxml import etree
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        fill_elem = shape.fill._fill
        solid = fill_elem.find(f'{{{ns}}}solidFill')
        if solid is not None:
            srgb = solid.find(f'{{{ns}}}srgbClr')
            if srgb is not None:
                alpha_elem = etree.SubElement(srgb, f'{{{ns}}}alpha')
                alpha_elem.set('val', str(int(alpha * 1000)))
    return shape


def _text_box(slide, left, top, width, height, text, font_size=18,
              color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullet_slide(slide, items, left, top, width, height,
                      font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT_GOLD):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(4)

        # Bullet character
        run_b = p.add_run()
        run_b.text = "▸ "
        run_b.font.size = Pt(font_size)
        run_b.font.color.rgb = bullet_color
        run_b.font.name = "Segoe UI"

        # Text
        run_t = p.add_run()
        run_t.text = item
        run_t.font.size = Pt(font_size)
        run_t.font.color.rgb = color
        run_t.font.name = "Segoe UI"
    return txBox


def _section_header(slide, number, title, subtitle=""):
    _add_rect(slide, Inches(0), Inches(0), W, Inches(1.6), RGBColor(0x15, 0x15, 0x28))
    # Number badge
    badge = _add_rect(slide, Inches(0.6), Inches(0.45), Inches(0.7), Inches(0.7), ACCENT_GOLD)
    badge_tf = badge.text_frame
    badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    badge_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = badge_tf.paragraphs[0].add_run()
    run.text = str(number)
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = BG_DARK
    run.font.name = "Segoe UI"

    _text_box(slide, Inches(1.6), Inches(0.35), Inches(10), Inches(0.6),
              title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        _text_box(slide, Inches(1.6), Inches(0.95), Inches(10), Inches(0.4),
                  subtitle, font_size=16, color=DIM_GRAY)


def _flow_boxes(slide, items, top, box_width=Inches(1.8), box_height=Inches(0.8), gap=Inches(0.15)):
    """Draw a horizontal flow: box → box → box ..."""
    total_w = len(items) * box_width + (len(items) - 1) * gap
    start_left = (W - total_w) // 2

    colors = [ACCENT_BLUE, ACCENT_CYAN, ACCENT_GOLD, ACCENT_PURPLE, ACCENT_ORANGE, ACCENT_RED, ACCENT_BLUE]
    for i, (label, accent) in enumerate(zip(items, colors)):
        left = start_left + i * (box_width + gap)
        box = _add_rect(slide, left, top, box_width, box_height, BG_CARD)
        # Top border line
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, box_width, Pt(4))
        border.fill.solid()
        border.fill.fore_color.rgb = accent
        border.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(13)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Segoe UI"

        # Arrow between boxes
        if i < len(items) - 1:
            arrow_left = left + box_width + Inches(0.01)
            _text_box(slide, arrow_left, top + Inches(0.15), gap, Inches(0.5),
                      "→", font_size=20, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_set_bg(slide)

# Gold accent bar at top
_add_rect(slide, Inches(0), Inches(0), W, Pt(5), ACCENT_GOLD)

# Title
_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
          "FilmDB", font_size=64, color=ACCENT_GOLD, bold=True,
          alignment=PP_ALIGN.CENTER, font_name="Segoe UI")

_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
          "Personalised Cinematic Intelligence", font_size=28, color=WHITE,
          alignment=PP_ALIGN.CENTER)

_text_box(slide, Inches(1), Inches(3.7), Inches(11), Inches(0.5),
          "Powered by Deterministic AI & Structured Knowledge Retrieval",
          font_size=16, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom bar
_add_rect(slide, Inches(4.5), Inches(5.0), Inches(4.3), Pt(2), ACCENT_GOLD)

_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
          "Project Overview — Architecture & Components",
          font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2: INTRODUCTION
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide)
_section_header(slide, "01", "Introduction", "What is FilmDB?")

_text_box(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.6),
          "FilmDB is a Q/A engine over a multi-layered movie knowledge base — capable of entity lookup, "
          "plot explanation, critic reviews, recommendations, filmography, and movie comparison.",
          font_size=18, color=LIGHT_GRAY)

# Key capabilities cards
caps = [
    ("🎬", "Entity Lookup", "Movie metadata, ratings,\ncast, and overview"),
    ("📖", "Plot Explanation", "Full Wikipedia plot\nsummaries with analysis"),
    ("✍️", "Critic Reviews", "Aggregated RT reviews\nwith sentiment analysis"),
    ("⭐", "Recommendations", "Tag-based similarity\nfrom 32M user ratings"),
    ("👤", "Filmography", "Director/actor career\nfrom 2.4M person index"),
    ("🔀", "Comparison", "Side-by-side movie\ndata comparison"),
]
for i, (icon, title, desc) in enumerate(caps):
    col = i % 3
    row = i // 3
    left = Inches(0.8) + col * Inches(4.1)
    top = Inches(3.0) + row * Inches(2.0)
    card = _add_rect(slide, left, top, Inches(3.8), Inches(1.7), BG_CARD)

    _text_box(slide, left + Inches(0.2), top + Inches(0.15), Inches(3.4), Inches(0.4),
              f"{icon}  {title}", font_size=18, color=ACCENT_GOLD, bold=True)
    _text_box(slide, left + Inches(0.2), top + Inches(0.65), Inches(3.4), Inches(0.9),
              desc, font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3: OVERALL ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide)
_section_header(slide, "02", "Overall Architecture", "End-to-end conversation pipeline")

# Flow diagram
flow_items = ["Intent\nAgent", "Entity\nResolver", "Routing\nMatrix",
              "Tool\nExecutor", "Query\nEngine", "Layout\nPolicy", "Response\nFormatter"]
_flow_boxes(slide, flow_items, top=Inches(2.2), box_width=Inches(1.5), box_height=Inches(0.9), gap=Inches(0.2))

_text_box(slide, Inches(0.8), Inches(1.8), Inches(3), Inches(0.3),
          "User Query →", font_size=14, color=ACCENT_CYAN, bold=True)
_text_box(slide, Inches(10.5), Inches(1.8), Inches(2.5), Inches(0.3),
          "→ Response", font_size=14, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.RIGHT)

# Description below
desc_items = [
    "No layer bypasses another — deterministic, traceable pipeline",
    "IntentAgent classifies user query into 23 intents via LLM",
    "EntityResolver maps titles to canonical IMDb IDs using KB + fuzzy matching",
    "RoutingMatrix selects which tools to call (KB-first, API fallback)",
    "ToolExecutor runs tools in parallel with caching and governance",
    "LayoutPolicy decides response format (FULL_CARD, RECOMMENDATION_GRID, ANALYSIS_TEXT…)",
]
_add_bullet_slide(slide, desc_items, Inches(0.8), Inches(3.6), Inches(11.5), Inches(3.5),
                  font_size=15, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4: INTENT AGENT
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide)
_section_header(slide, "03", "Layer 1 — Intent Agent", "LLM-powered intent classification")

_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.5),
          "What it does:", font_size=20, color=ACCENT_GOLD, bold=True)
_text_box(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(1.2),
          "Classifies the user's message into one of 23 intents, extracts entities "
          "(movie titles, person names, genres, years), and assigns a confidence score.",
          font_size=15, color=LIGHT_GRAY)

_text_box(slide, Inches(0.8), Inches(3.8), Inches(5.5), Inches(0.4),
          "23 Supported Intents:", font_size=16, color=ACCENT_CYAN, bold=True)

intents_left = [
    "ENTITY_LOOKUP", "PLOT_EXPLANATION", "CRITIC_SUMMARY",
    "MOVIE_SIMILARITY", "RECOMMENDATION", "TOP_RATED",
    "FILMOGRAPHY", "PERSON_LOOKUP", "COMPARISON",
    "STREAMING_AVAILABILITY", "TRENDING", "UPCOMING",
]
intents_right = [
    "REVIEWS", "AVAILABILITY", "DOWNLOAD",
    "LEGAL_DOWNLOAD", "ILLEGAL_DOWNLOAD_REQUEST",
    "STREAMING_DISCOVERY", "AWARD_LOOKUP",
    "ANALYTICAL_EXPLANATION", "OFFICIAL",
    "GREETING", "GENERAL_CONVERSATION",
]

_add_bullet_slide(slide, intents_left, Inches(0.8), Inches(4.3), Inches(5), Inches(3),
                  font_size=13, color=LIGHT_GRAY, bullet_color=ACCENT_BLUE)
_add_bullet_slide(slide, intents_right, Inches(6.5), Inches(4.3), Inches(5), Inches(3),
                  font_size=13, color=LIGHT_GRAY, bullet_color=ACCENT_BLUE)

# Example box
_add_rect(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(1.8), BG_CARD)
_text_box(slide, Inches(7.2), Inches(2.1), Inches(5), Inches(0.3),
          "Example", font_size=14, color=ACCENT_GOLD, bold=True)
_text_box(slide, Inches(7.2), Inches(2.5), Inches(5), Inches(1.2),
          '"movies like Inception"\n→ MOVIE_SIMILARITY\n→ entities: [{type: movie, value: "Inception"}]\n→ confidence: 90',
          font_size=13, color=LIGHT_GRAY, font_name="Consolas")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5: ENTITY RESOLVER
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide)
_section_header(slide, "04", "Layer 2 — Entity Resolver", "Title normalization & IMDb ID resolution")

_text_box(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.5),
          "What it does:", font_size=20, color=ACCENT_GOLD, bold=True)
items = [
    "Strips filler phrases from user input (\"tell me about\", \"where can I watch\"…)",
    "Resolves aliases: \"Dark Knight\" → \"The Dark Knight\", \"Shawshank\" → \"The Shawshank Redemption\"",
    "Infers entity type: movie, person, award_event, catalog",
    "Resolves movie titles to IMDb ID (tconst) via KB with fuzzy matching (≥90% threshold)",
    "Extracts year from message for disambiguation",
    "Detects public domain status (pre-1928 films)",
    "Attaches canonical_id (imdb_id) so all downstream tools get it for free",
]
_add_bullet_slide(slide, items, Inches(0.8), Inches(2.6), Inches(11.5), Inches(4.5),
                  font_size=16, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6: ROUTING MATRIX
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide)
_section_header(slide, "05", "Layer 3 — Routing Matrix", "Intent → Tool mapping with KB-first strategy")

_text_box(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.5),
          "What it does:", font_size=20, color=ACCENT_GOLD, bold=True)
_text_box(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(0.8),
          "Maps each intent to required/optional/forbidden tools. Uses KB-first routing — "
          "local Parquet data is primary, external APIs are fallback.",
          font_size=15, color=LIGHT_GRAY)

# Example routing table
routes = [
    ("ENTITY_LOOKUP",  "kb_entity",     "imdb, wikipedia, watchmode"),
    ("PLOT_EXPLANATION","kb_plot",       "kb_entity, wikipedia"),
    ("CRITIC_SUMMARY", "kb_critic",     "rt_reviews, web_search"),
    ("MOVIE_SIMILARITY","kb_similarity", "similarity (API)"),
    ("TOP_RATED",      "kb_top_rated",  "imdb_top_rated_english"),
    ("FILMOGRAPHY",    "kb_filmography","imdb_person, wikipedia"),
    ("COMPARISON",     "kb_comparison", "web_search"),
    ("STREAMING",      "watchmode",     "web_search"),
    ("TRENDING",       "web_search",    "imdb_trending_tamil"),
]

# Table header
_add_rect(slide, Inches(0.8), Inches(3.5), Inches(3.5), Inches(0.45), ACCENT_GOLD)
_add_rect(slide, Inches(4.3), Inches(3.5), Inches(3.5), Inches(0.45), ACCENT_GOLD)
_add_rect(slide, Inches(7.8), Inches(3.5), Inches(4.7), Inches(0.45), ACCENT_GOLD)
_text_box(slide, Inches(0.9), Inches(3.52), Inches(3.3), Inches(0.4),
          "Intent", font_size=14, color=BG_DARK, bold=True)
_text_box(slide, Inches(4.4), Inches(3.52), Inches(3.3), Inches(0.4),
          "Required (KB-first)", font_size=14, color=BG_DARK, bold=True)
_text_box(slide, Inches(7.9), Inches(3.52), Inches(4.5), Inches(0.4),
          "Optional (Fallback)", font_size=14, color=BG_DARK, bold=True)

for i, (intent, req, opt) in enumerate(routes):
    y = Inches(4.0) + i * Inches(0.38)
    bg = BG_CARD if i % 2 == 0 else RGBColor(0x14, 0x14, 0x24)
    _add_rect(slide, Inches(0.8), y, Inches(11.7), Inches(0.38), bg)
    _text_box(slide, Inches(0.9), y + Pt(2), Inches(3.3), Inches(0.35),
              intent, font_size=12, color=ACCENT_CYAN, font_name="Consolas")
    _text_box(slide, Inches(4.4), y + Pt(2), Inches(3.3), Inches(0.35),
              req, font_size=12, color=WHITE)
    _text_box(slide, Inches(7.9), y + Pt(2), Inches(4.5), Inches(0.35),
              opt, font_size=12, color=DIM_GRAY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7: TOOL EXECUTOR & QUERY ENGINE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide)
_section_header(slide, "06", "Layer 4 & 5 — Tool Executor & Query Engine",
                "Parallel execution with caching + KB data access")

# Left side: Tool Executor
_add_rect(slide, Inches(0.6), Inches(2.0), Inches(5.8), Inches(5.2), BG_CARD)
_text_box(slide, Inches(0.8), Inches(2.1), Inches(5.4), Inches(0.4),
          "Tool Executor", font_size=22, color=ACCENT_GOLD, bold=True)
items_left = [
    "Executes tools in parallel (asyncio.gather)",
    "Built-in caching: metadata, streaming, similarity",
    "Governance: validates tool calls, max 4 tools/request",
    "Handles 18 tools (11 API + 7 KB)",
    "Stores tool call traces for debugging",
    "Cache hits skip API calls entirely",
]
_add_bullet_slide(slide, items_left, Inches(0.8), Inches(2.7), Inches(5.4), Inches(4.2),
                  font_size=14, color=LIGHT_GRAY)

# Right side: Query Engine
_add_rect(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(5.2), BG_CARD)
_text_box(slide, Inches(7.0), Inches(2.1), Inches(5.4), Inches(0.4),
          "FilmDB Query Engine", font_size=22, color=ACCENT_CYAN, bold=True)
items_right = [
    "Central data access layer for all KB tools",
    "Loads 7 Parquet layers once at startup (~400 MB)",
    "Title → IMDb ID resolution (exact + fuzzy)",
    "Provides: entity_lookup(), plot_analysis()",
    "   critic_summary(), movie_similarity()",
    "   top_rated(), person_filmography()",
    "   compare_movies()",
]
_add_bullet_slide(slide, items_right, Inches(7.0), Inches(2.7), Inches(5.4), Inches(4.2),
                  font_size=14, color=LIGHT_GRAY, bullet_color=ACCENT_CYAN)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8: LAYOUT POLICY & RESPONSE FORMATTER
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide)
_section_header(slide, "07", "Layer 6 & 7 — Layout Policy & Response Formatter",
                "Structured output with LLM synthesis")

_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.4),
          "Layout Policy — Response Modes:", font_size=20, color=ACCENT_GOLD, bold=True)
modes = [
    ("FULL_CARD", "Complete movie card with poster, rating, streaming"),
    ("MINIMAL_CARD", "Lightweight card for download/legal queries"),
    ("RECOMMENDATION_GRID", "Grid of movie recommendations"),
    ("EXPLANATION_ONLY", "Text-only analytical response"),
    ("ANALYSIS_TEXT", "Plot/critic analysis (new)"),
    ("FILMOGRAPHY_LIST", "Person career listing (new)"),
    ("COMPARISON_TABLE", "Side-by-side movie comparison (new)"),
    ("AVAILABILITY_FOCUS", "Streaming platform availability"),
]
_add_bullet_slide(slide, [f"{m}: {d}" for m, d in modes],
                  Inches(0.8), Inches(2.6), Inches(11.5), Inches(4.5),
                  font_size=14, color=LIGHT_GRAY)

_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
          "Response Formatter: Summarizes tool outputs → builds LLM prompt → generates natural language response",
          font_size=14, color=DIM_GRAY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9: DATASETS
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide)
_section_header(slide, "08", "Datasets & Knowledge Base", "Multi-source movie knowledge graph")

datasets = [
    ("IMDb",            "title.basics, ratings,\nprincipals, crew, names",  "12.3M titles\n730K movies",  "~5 GB",    ACCENT_BLUE),
    ("TMDB",            "Movie metadata v11\n(overview, poster, budget)",   "1.4M movies\n347K matched",  "1.2 GB",   ACCENT_CYAN),
    ("MovieLens 32M",   "User ratings + tags\n(GroupLens research)",        "32M ratings\n87K movies",    "1.1 GB",   ACCENT_GOLD),
    ("Rotten Tomatoes",  "Critic reviews +\nsentiment scores",              "1M reviews\n143K movies",    "1 GB",     ACCENT_ORANGE),
    ("Wikipedia Plots",  "Full movie plot\nsummaries",                      "35K plots\n29.7K matched",   "86 MB",    ACCENT_PURPLE),
    ("Indian Movies",    "Regional cinema\n(Hindi, Tamil, etc.)",           "50.6K films\n7 languages",   "9.5 MB",   ACCENT_RED),
]
for i, (name, desc, stats, size, color) in enumerate(datasets):
    col = i % 3
    row = i // 3
    left = Inches(0.6) + col * Inches(4.15)
    top = Inches(2.0) + row * Inches(2.6)
    card = _add_rect(slide, left, top, Inches(3.9), Inches(2.3), BG_CARD)
    # Color strip
    _add_rect(slide, left, top, Inches(3.9), Pt(4), color)

    _text_box(slide, left + Inches(0.15), top + Inches(0.2), Inches(3.6), Inches(0.4),
              name, font_size=18, color=color, bold=True)
    _text_box(slide, left + Inches(0.15), top + Inches(0.65), Inches(2.2), Inches(0.8),
              desc, font_size=12, color=LIGHT_GRAY)
    _text_box(slide, left + Inches(2.4), top + Inches(0.65), Inches(1.4), Inches(0.8),
              stats, font_size=12, color=WHITE, alignment=PP_ALIGN.RIGHT)
    _text_box(slide, left + Inches(0.15), top + Inches(1.7), Inches(3.6), Inches(0.4),
              f"Raw Size: {size}", font_size=11, color=DIM_GRAY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10: KB MAPPING PIPELINE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide)
_section_header(slide, "09", "KB Mapping Pipeline", "Raw datasets → 7 canonical Parquet layers")

layers = [
    ("movie_entity.parquet",         "730K rows",   "21.5 MB",  "Master movie table (IMDb ID primary key)"),
    ("person_index.parquet",         "2.37M rows",  "162.3 MB", "Pre-aggregated person filmography"),
    ("metadata_layer.parquet",       "347K rows",   "87.8 MB",  "TMDB enrichment (overview, poster…)"),
    ("plot_layer.parquet",           "29.7K rows",  "39.6 MB",  "Wikipedia full plot summaries"),
    ("review_layer.parquet",         "927K rows",   "80.5 MB",  "Rotten Tomatoes critic reviews"),
    ("recommendation_layer.parquet", "87.6K rows",  "4.5 MB",   "MovieLens aggregated ratings/tags"),
    ("regional_layer.parquet",       "50.6K rows",  "1.3 MB",   "Indian cinema regional data"),
]

# Table header
_add_rect(slide, Inches(0.6), Inches(2.1), Inches(4.2), Inches(0.45), ACCENT_GOLD)
_add_rect(slide, Inches(4.8), Inches(2.1), Inches(1.5), Inches(0.45), ACCENT_GOLD)
_add_rect(slide, Inches(6.3), Inches(2.1), Inches(1.2), Inches(0.45), ACCENT_GOLD)
_add_rect(slide, Inches(7.5), Inches(2.1), Inches(5.2), Inches(0.45), ACCENT_GOLD)
_text_box(slide, Inches(0.7), Inches(2.12), Inches(4), Inches(0.4),
          "Output Layer", font_size=14, color=BG_DARK, bold=True)
_text_box(slide, Inches(4.9), Inches(2.12), Inches(1.3), Inches(0.4),
          "Rows", font_size=14, color=BG_DARK, bold=True)
_text_box(slide, Inches(6.4), Inches(2.12), Inches(1), Inches(0.4),
          "Size", font_size=14, color=BG_DARK, bold=True)
_text_box(slide, Inches(7.6), Inches(2.12), Inches(5), Inches(0.4),
          "Description", font_size=14, color=BG_DARK, bold=True)

for i, (name, rows, size, desc) in enumerate(layers):
    y = Inches(2.6) + i * Inches(0.5)
    bg = BG_CARD if i % 2 == 0 else RGBColor(0x14, 0x14, 0x24)
    _add_rect(slide, Inches(0.6), y, Inches(12.1), Inches(0.5), bg)
    _text_box(slide, Inches(0.7), y + Pt(3), Inches(4), Inches(0.4),
              name, font_size=13, color=ACCENT_CYAN, font_name="Consolas")
    _text_box(slide, Inches(4.9), y + Pt(3), Inches(1.3), Inches(0.4),
              rows, font_size=13, color=WHITE)
    _text_box(slide, Inches(6.4), y + Pt(3), Inches(1), Inches(0.4),
              size, font_size=13, color=WHITE)
    _text_box(slide, Inches(7.6), y + Pt(3), Inches(5), Inches(0.4),
              desc, font_size=13, color=LIGHT_GRAY)

_text_box(slide, Inches(0.6), Inches(6.3), Inches(12), Inches(0.5),
          "Total output: ~400 MB Parquet  |  Pipeline runtime: ~14 minutes  |  IMDb ID (tconst) as global key",
          font_size=14, color=DIM_GRAY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11: APIs & TOOLS
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide)
_section_header(slide, "10", "APIs & Tools", "18 tools — 7 KB-powered + 11 API-based")

# KB Tools
_text_box(slide, Inches(0.6), Inches(2.0), Inches(6), Inches(0.4),
          "KB Tools (Local, Instant)", font_size=20, color=ACCENT_CYAN, bold=True)
kb_tools = [
    "kb_entity — Movie metadata from Parquet KB",
    "kb_plot — Wikipedia plot text retrieval",
    "kb_critic — RT review aggregation + sentiment",
    "kb_similarity — MovieLens tag-based similarity",
    "kb_top_rated — Genre/year/language filtered ranking",
    "kb_filmography — Person career from 2.4M person index",
    "kb_comparison — Side-by-side two-movie data",
]
_add_bullet_slide(slide, kb_tools, Inches(0.6), Inches(2.5), Inches(6), Inches(3.5),
                  font_size=13, color=LIGHT_GRAY, bullet_color=ACCENT_CYAN)

# API Tools
_text_box(slide, Inches(6.8), Inches(2.0), Inches(6), Inches(0.4),
          "API Tools (Live Data)", font_size=20, color=ACCENT_ORANGE, bold=True)
api_tools = [
    "imdb (TMDB) — Live movie metadata + poster",
    "wikipedia — Live Wikipedia summaries",
    "watchmode — Streaming availability",
    "similarity — TMDB-based recommendations",
    "web_search — Serper API for live search",
    "rt_reviews — Rotten Tomatoes live reviews",
    "imdb_person — TMDB person lookup",
    "imdb_trending_tamil — Trending Tamil movies",
    "imdb_top_rated_english — Top rated English",
    "imdb_upcoming — Upcoming releases by region",
    "archive — Internet Archive public domain",
]
_add_bullet_slide(slide, api_tools, Inches(6.8), Inches(2.5), Inches(6), Inches(5),
                  font_size=13, color=LIGHT_GRAY, bullet_color=ACCENT_ORANGE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12: CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide)

_add_rect(slide, Inches(0), Inches(0), W, Pt(5), ACCENT_GOLD)

_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
          "Summary", font_size=40, color=ACCENT_GOLD, bold=True,
          alignment=PP_ALIGN.CENTER)

summary_items = [
    "Deterministic pipeline: Intent → Entity → Route → Execute → Format",
    "KB-first architecture: 7 Parquet layers (400 MB) for instant local data access",
    "18 tools total: 7 KB-powered (instant) + 11 API-based (live fallback)",
    "730K movies, 2.4M people, 927K reviews — all keyed by IMDb ID",
    "Extensible: analysis_layer placeholder ready for web-scraped articles",
    "Traceable: every request produces a full debug trace",
]
_add_bullet_slide(slide, summary_items, Inches(1.5), Inches(2.8), Inches(10), Inches(3.5),
                  font_size=18, color=LIGHT_GRAY, bullet_color=ACCENT_GOLD)

_add_rect(slide, Inches(4.5), Inches(6.2), Inches(4.3), Pt(2), ACCENT_GOLD)
_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
          "FilmDB — Personalised Cinematic Intelligence",
          font_size=16, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)


# ─── Save ──────────────────────────────────────────────────────────────────
output_path = r"D:\Evolve_Robot_Lab\Project\FilmDB_Demo\FilmDB_Overview.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
